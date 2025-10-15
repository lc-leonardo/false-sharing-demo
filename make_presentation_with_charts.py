import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont

# Data (from user-run)
runs = [
    (76.00, 64.00, 1.19),
    (71.00, 81.00, 0.88),
    (78.00, 74.00, 1.05),
]
averages = (75.00, 73.00)
throughput = (2666666667, 2739726027)

out_dir = '.'
# 1) Create timing chart (per-run)
run_ids = [1,2,3]
fs_times = [r[0] for r in runs]
opt_times = [r[1] for r in runs]

plt.figure(figsize=(8,4))
plt.plot(run_ids, fs_times, marker='o', label='False sharing')
plt.plot(run_ids, opt_times, marker='o', label='Optimized (padded)')
plt.xticks(run_ids)
plt.xlabel('Run')
plt.ylabel('Execution time (ms)')
plt.title('Per-run Execution Time')
plt.legend()
plt.grid(True, linestyle='--', alpha=0.5)
plt.tight_layout()
plt.savefig(f'{out_dir}/timing_per_run.png', dpi=150)
plt.close()

# 2) Create throughput bar chart (averages)
labels = ['False sharing', 'Optimized']
vals = list(throughput)
vals_mill = [v/1e6 for v in vals]

plt.figure(figsize=(6,4))
bars = plt.bar(labels, vals_mill, color=['tab:orange','tab:blue'])
plt.ylabel('Throughput (million ops/sec)')
plt.title('Average Throughput')
for bar, val in zip(bars, vals_mill):
    plt.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 10, f'{val:.2f}', ha='center')
plt.tight_layout()
plt.savefig(f'{out_dir}/throughput_avg.png', dpi=150)
plt.close()

# 3) Create illustrative diagrams: false sharing vs padded
W, H = 800, 200
font = None
try:
    font = ImageFont.truetype('arial.ttf', 14)
except Exception:
    font = ImageFont.load_default()

# False sharing image
img = Image.new('RGB', (W,H), 'white')
d = ImageDraw.Draw(img)
# draw cache line
d.rectangle([50,50,750,150], outline='black', width=2)
# draw counters close together
x = 70
colors = ['red','green','blue','purple']
for i in range(4):
    d.rectangle([x,70,x+60,130], fill=colors[i], outline='black')
    d.text((x+5,135), f'counter[{i}]', fill='black', font=font)
    x += 80

d.text((50,20), 'False sharing: many counters share one cache line (invalidations)', fill='black', font=font)
img.save(f'{out_dir}/illustration_false_sharing.png')

# Padded image
img2 = Image.new('RGB', (W,H), 'white')
d2 = ImageDraw.Draw(img2)
# Draw 4 separate cache line blocks
x = 50
for i in range(4):
    d2.rectangle([x,40,x+160,160], outline='black', width=2)
    d2.rectangle([x+10,60,x+70,120], fill=colors[i], outline='black')
    d2.text((x+10,125), f'counter[{i}]', fill='black', font=font)
    x += 180

d2.text((50,10), 'Padded: each counter in its own cache line → no ping-pong', fill='black', font=font)
img2.save(f'{out_dir}/illustration_padded.png')

# 4) Build PPTX embedding images
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Title slide
sld = prs.slides.add_slide(prs.slide_layouts[0])
sld.shapes.title.text = 'False Sharing — Demo & Mitigation'
sld.placeholders[1].text = 'Timing and memory-layout charts included'

# Timing slide
sld = prs.slides.add_slide(prs.slide_layouts[5])
sld.shapes.title.text = 'Per-run Timing'
left = Inches(1)
top = Inches(1.2)
pic = sld.shapes.add_picture('timing_per_run.png', left, top, width=Inches(10))

# Throughput slide
sld = prs.slides.add_slide(prs.slide_layouts[5])
sld.shapes.title.text = 'Average Throughput'
left = Inches(1)
top = Inches(1.2)
pic = sld.shapes.add_picture('throughput_avg.png', left, top, width=Inches(8))

# Illustrations slide
sld = prs.slides.add_slide(prs.slide_layouts[5])
sld.shapes.title.text = 'Illustration: False Sharing vs Padded'
left = Inches(0.5)
top = Inches(1.0)
pic1 = sld.shapes.add_picture('illustration_false_sharing.png', left, top, width=Inches(6))
pic2 = sld.shapes.add_picture('illustration_padded.png', Inches(7), top, width=Inches(6))

# Add a slide with summary bullet points
sld = prs.slides.add_slide(prs.slide_layouts[1])
sld.shapes.title.text = 'Summary & Recommendations'
body = sld.shapes.placeholders[1].text_frame
body.text = 'Results summary:'
body.add_paragraph().text = f"Avg speedup (padding): {averages[0]/averages[1]:.2f}x (~{(averages[0]-averages[1]) / averages[0] * 100:.1f}% reduction)"
body.add_paragraph().text = 'Recommendations: use padding, thread pinning, and larger workloads'

out_pptx = 'false_sharing_presentation_with_charts.pptx'
prs.save(out_pptx)
print('Wrote', out_pptx)
