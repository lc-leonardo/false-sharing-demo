from pptx import Presentation
from pptx.util import Inches, Pt

# Data from user's run
system_info = {
    'processors': 16,
    'max_threads': 16,
    'test_threads': 4,
    'iterations_per_thread': 50000000,
    'total_ops': 200000000,
    'cache_line': 64,
}

memory_layout = {
    'unoptimized': {
        'sizeof': '16 bytes',
        'spacing': '4 bytes',
        'offset0': '0 bytes',
        'offset1': '4 bytes',
        'note': 'Multiple counters in same cache line'
    },
    'optimized': {
        'sizeof': '64 bytes',
        'spacing': '64 bytes',
        'offset0': '0 bytes',
        'offset1': '64 bytes',
        'note': 'Each counter in separate cache line'
    }
}

runs = [
    (76.00, 64.00, 1.19),
    (71.00, 81.00, 0.88),
    (78.00, 74.00, 1.05),
]

averages = {
    'false_sharing': 75.00,
    'optimized': 73.00,
    'speedup': 1.03,
}

throughput = {
    'false_sharing': 2666666667,
    'optimized': 2739726027,
    'gain': 1.03,
}

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Helper to add a title slide
def add_title_slide(title, subtitle=None):
    sld = prs.slides.add_slide(prs.slide_layouts[0])
    title_tf = sld.shapes.title
    title_tf.text = title
    if subtitle:
        subtitle_tf = sld.placeholders[1]
        subtitle_tf.text = subtitle

# Helper: add a simple title + bullet slide
def add_bullet_slide(title, bullets):
    sld = prs.slides.add_slide(prs.slide_layouts[1])
    sld.shapes.title.text = title
    body = sld.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
            p.text = b
        else:
            p = body.add_paragraph()
            p.text = b
        p.level = 0
        p.font.size = Pt(18)

# Build slides
add_title_slide('False Sharing — Demo & Mitigation', 'Dubois et al. (1990) & Jeremiassen & Eggers (1993)')

add_bullet_slide('Agenda', [
    'What is false sharing?',
    'Demonstration code and mitigation approach',
    'Memory-layout evidence',
    'Measured results (your run)',
    'Analysis and recommendations',
])

add_bullet_slide('What is False Sharing?', [
    'Independent variables used by different threads share a cache line',
    'Writes cause cache-line invalidation and coherence traffic',
    'Leads to significant slowdowns in multithreaded code',
])

add_bullet_slide('Unoptimized Code (concept)', [
    'Shared array of counters packed contiguously',
    'Each thread increments counter[tid] rapidly',
    'Counters fall into same cache line → false sharing',
])

add_bullet_slide('Optimized Code (compile-time padding)', [
    'Per-thread padded struct: int counter; char pad[64 - sizeof(int)];',
    'Each counter aligned to its own cache line (alignas(64) or __attribute__)',
    'Eliminates cache invalidation between threads',
])

# Memory layout slide
mem_bullets = [
    'Unoptimized (False Sharing):',
    f"  sizeof(SharedData): {memory_layout['unoptimized']['sizeof']}",
    f"  Counter spacing: {memory_layout['unoptimized']['spacing']}",
    f"  counter[0] at offset: {memory_layout['unoptimized']['offset0']}",
    f"  counter[1] at offset: {memory_layout['unoptimized']['offset1']}",
    f"  → {memory_layout['unoptimized']['note']}",
    '',
    'Optimized (Padded):',
    f"  sizeof(PaddedCounter): {memory_layout['optimized']['sizeof']}",
    f"  Counter spacing: {memory_layout['optimized']['spacing']}",
    f"  padded_counter[0] at offset: {memory_layout['optimized']['offset0']}",
    f"  padded_counter[1] at offset: {memory_layout['optimized']['offset1']}",
    f"  → {memory_layout['optimized']['note']}",
]
add_bullet_slide('Memory Layout — Measured', mem_bullets)

# System info slide
sys_bullets = [
    f"Processors: {system_info['processors']}",
    f"Max threads: {system_info['max_threads']}",
    f"Test threads: {system_info['test_threads']}",
    f"Iterations/thread: {system_info['iterations_per_thread']}",
    f"Total ops/run: {system_info['total_ops']}",
    f"Cache line size: {system_info['cache_line']} bytes",
]
add_bullet_slide('System & Test Parameters', sys_bullets)

# Results slide (per-run)
run_lines = []
for idx, (fs, opt, su) in enumerate(runs, start=1):
    run_lines.append(f"Run {idx}: False sharing {fs:.2f} ms, Optimized {opt:.2f} ms, Speedup {su:.2f}x")
add_bullet_slide('Per-Run Results', run_lines)

# Final averages slide
add_bullet_slide('Final Averages', [
    f"Avg without padding: {averages['false_sharing']:.2f} ms",
    f"Avg with padding:    {averages['optimized']:.2f} ms",
    f"Speedup:             {averages['speedup']:.2f}x ({(averages['speedup'] - 1.0) * 100.0:.1f}% time reduction)",
    '',
    f"Throughput (False sharing): {throughput['false_sharing']:,} ops/sec",
    f"Throughput (Optimized):     {throughput['optimized']:,} ops/sec",
    f"Throughput gain:            {throughput['gain']:.2f}x",
])

add_bullet_slide('Analysis', [
    'Padding reduced false sharing; observed improvement was small (≈1.03x)',
    'Possible reasons: hardware cache-coherence efficiency, measurement noise',
    'Recommendations: increase workload, pin threads, measure cache misses',
])

add_bullet_slide('Practical Recommendations', [
    'Use padding / alignas(64) for per-thread hot data',
    'Prefer compile-time layout changes when access patterns are static',
    'Combine with thread pinning and larger workloads for clearer gains',
])

add_bullet_slide('References', [
    'Dubois, Scheurich & Briggs (1990) - False sharing analysis',
    'Jeremiassen & Eggers (1993) - Compile-time data transformations',
])

# Save presentation
out_path = 'false_sharing_presentation.pptx'
prs.save(out_path)
print('Wrote', out_path)
